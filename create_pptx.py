"""Generate a PromptKit intro presentation (5–10 slides, engineering audience)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme colours ──────────────────────────────────────────────
BG        = RGBColor(0x1E, 0x1E, 0x2E)   # dark background
ACCENT    = RGBColor(0x89, 0xB4, 0xFA)   # soft blue
ACCENT2   = RGBColor(0xA6, 0xE3, 0xA1)   # green
ACCENT3   = RGBColor(0xF9, 0xE2, 0xAF)   # yellow/gold
WHITE     = RGBColor(0xCD, 0xD6, 0xF4)   # off-white text
DIM       = RGBColor(0x6C, 0x70, 0x86)   # muted text
SURFACE   = RGBColor(0x31, 0x32, 0x44)   # card/box bg
CODE_BG   = RGBColor(0x18, 0x18, 0x25)   # code block bg
ORANGE    = RGBColor(0xFA, 0xB3, 0x87)

W = Inches(13.333)  # 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── helpers ────────────────────────────────────────────────────
def dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_text(slide, left, top, width, height, text, *,
             font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Segoe UI", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(font_size * 0.3)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf

def add_para(tf, text, *, font_size=18, color=WHITE, bold=False,
             font_name="Segoe UI", alignment=PP_ALIGN.LEFT, space_after=None,
             line_spacing=None):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_after is not None:
        p.space_after = Pt(space_after)
    if line_spacing is not None:
        p.line_spacing = Pt(line_spacing)
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color=SURFACE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_accent_bar(slide, left, top, width=Inches(0.06), height=Inches(0.5),
                   color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_bullet_block(slide, left, top, width, bullets, *,
                     font_size=16, color=WHITE, bullet_color=ACCENT,
                     heading=None, heading_color=ACCENT):
    h = Inches(0.4) if heading else Inches(0)
    total_h = h + Inches(len(bullets) * 0.42)
    tf = add_text(slide, left, top, width, total_h,
                  heading or bullets[0],
                  font_size=font_size + 2 if heading else font_size,
                  color=heading_color if heading else color,
                  bold=bool(heading))
    start = 0 if heading else 1
    for b in bullets[start:]:
        add_para(tf, f"  •  {b}", font_size=font_size, color=color,
                 space_after=4, line_spacing=font_size * 1.3)
    return tf

def title_strip(slide, title, subtitle=None):
    """Top accent bar + title across the slide."""
    # accent line
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.06)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = ACCENT
    slide.shapes[-1].line.fill.background()
    # title
    add_text(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.65),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.45),
                 subtitle, font_size=16, color=DIM)

def speaker_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
dark_bg(s)

# Large accent line
s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.08)
                   ).fill.solid()
s.shapes[-1].fill.fore_color.rgb = ACCENT
s.shapes[-1].line.fill.background()

# Logo
logo_size = Inches(2.5)
s.shapes.add_picture("PromptKit-logo.png",
                     Inches(0.9), Inches(1.2), logo_size, logo_size)

# Title text — shifted right to sit beside the logo
add_text(s, Inches(3.7), Inches(1.6), Inches(9), Inches(1.2),
         "PromptKit", font_size=64, color=WHITE, bold=True)
add_text(s, Inches(3.7), Inches(2.8), Inches(9), Inches(0.8),
         "Composable, version-controlled prompts\nfor AI-assisted software engineering",
         font_size=24, color=ACCENT)
add_text(s, Inches(3.7), Inches(4.1), Inches(9), Inches(0.5),
         "github.com/microsoft/promptkit  ·  MIT License  ·  v0.2.0",
         font_size=14, color=DIM)
add_text(s, Inches(0.9), Inches(6.2), Inches(11), Inches(0.4),
         "Engineering Deep-Dive  ·  Getting Started",
         font_size=14, color=DIM)
speaker_notes(s,
    "Welcome slide. PromptKit treats prompts as engineered software components — "
    "composable, testable, version-controlled. Today we'll cover the architecture, "
    "what's included, and how to start using it immediately.")


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "The Problem")

problems = [
    ("Ad-hoc prompts", "Written once, pasted everywhere, never improved"),
    ("No version control", "Can't track what changed, who changed it, or why"),
    ("No reuse", "Every engineer reinvents the same persona / guardrails / format"),
    ("No testing", "No way to verify prompt quality or catch regressions"),
    ("Inconsistent outputs", "Same task → wildly different results across the team"),
]

for i, (title, desc) in enumerate(problems):
    y = Inches(1.65 + i * 1.0)
    add_accent_bar(s, Inches(0.9), y, height=Inches(0.6),
                   color=[ORANGE, ACCENT3, ACCENT, ACCENT2, ACCENT][i])
    add_text(s, Inches(1.15), y - Inches(0.03), Inches(4), Inches(0.35),
             title, font_size=20, color=WHITE, bold=True)
    add_text(s, Inches(1.15), y + Inches(0.32), Inches(10), Inches(0.35),
             desc, font_size=15, color=DIM)

speaker_notes(s,
    "Walk through each pain point. Most teams treat prompts as throwaway text. "
    "But prompts are the primary interface to AI — they deserve the same engineering "
    "rigor we apply to code: modularity, review, testing, reuse.")


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — The Fix: Prompts as Code
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "The Fix: Treat Prompts as Code")

pillars = [
    ("Modular", "Compose from reusable\nbuilding blocks", ACCENT),
    ("Version-Controlled", "Track changes via\nGit like any code", ACCENT2),
    ("Testable", "Reference-compare\nagainst known-good", ACCENT3),
    ("Extensible", "Add your own personas,\nprotocols, templates", ORANGE),
]

card_w = Inches(2.6)
gap = Inches(0.35)
total = len(pillars) * card_w.inches + (len(pillars) - 1) * gap.inches
start_x = (13.333 - total) / 2

for i, (title, desc, color) in enumerate(pillars):
    x = Inches(start_x + i * (card_w.inches + gap.inches))
    y = Inches(2.2)
    add_rounded_rect(s, x, y, card_w, Inches(2.8), SURFACE)
    # color bar at top of card
    s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.06)
                       ).fill.solid()
    s.shapes[-1].fill.fore_color.rgb = color
    s.shapes[-1].line.fill.background()
    add_text(s, x + Inches(0.25), y + Inches(0.35), Inches(2.1), Inches(0.4),
             title, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.25), y + Inches(0.95), Inches(2.1), Inches(1.2),
             desc, font_size=15, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text(s, Inches(1.5), Inches(5.5), Inches(10), Inches(0.8),
         "PromptKit applies the same engineering discipline you use for software\n"
         "to the prompts that build it.",
         font_size=16, color=DIM, alignment=PP_ALIGN.CENTER)

speaker_notes(s,
    "Four engineering principles applied to prompts. Modularity means you compose "
    "from reusable layers instead of writing monolithic prompts. Version control means "
    "prompts live in Git with full history. Testable means we have a reference-comparison "
    "methodology. Extensible means the team can contribute new components via PR.")


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture: The 5-Layer Stack
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Architecture: The 5-Layer Composition Stack")

layers = [
    ("① Persona",    "Who the LLM becomes",              "systems-engineer, security-auditor, devops-engineer …",  ACCENT),
    ("② Protocols",  "How it reasons + guardrails",       "anti-hallucination, root-cause-analysis, thread-safety …", ACCENT2),
    ("③ Format",     "Output structure & rules",          "requirements-doc, investigation-report, pipeline-spec …",   ACCENT3),
    ("④ Taxonomy",   "Domain classification (optional)",  "Severity levels, priority tiers, risk categories",          ORANGE),
    ("⑤ Template",   "The task with parameters",          "investigate-bug, review-code, author-design-doc …",         RGBColor(0xF3, 0x8B, 0xA8)),
]

for i, (name, role, examples, color) in enumerate(layers):
    y = Inches(1.55 + i * 1.05)
    # Layer bar
    bar_w = Inches(11.5 - i * 0.6)
    x = Inches((13.333 - bar_w.inches) / 2)
    add_rounded_rect(s, x, y, bar_w, Inches(0.85), SURFACE)
    # color accent on left
    s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.07), Inches(0.85)
                       ).fill.solid()
    s.shapes[-1].fill.fore_color.rgb = color
    s.shapes[-1].line.fill.background()
    # text
    add_text(s, x + Inches(0.25), y + Inches(0.05), Inches(2.2), Inches(0.4),
             name, font_size=18, color=color, bold=True)
    add_text(s, x + Inches(2.5), y + Inches(0.05), Inches(3), Inches(0.4),
             role, font_size=15, color=WHITE)
    add_text(s, x + Inches(2.5), y + Inches(0.42), bar_w - Inches(2.7), Inches(0.4),
             examples, font_size=12, color=DIM)

add_text(s, Inches(1), Inches(6.9), Inches(11), Inches(0.4),
         "Templates declare which components to compose via YAML frontmatter → bootstrap engine snaps them together",
         font_size=13, color=DIM, alignment=PP_ALIGN.CENTER)

speaker_notes(s,
    "Walk through the stack bottom-up: Template is the task you want to accomplish. "
    "It declares which persona, protocols, and format to use. The bootstrap engine "
    "reads these declarations and assembles a single coherent prompt in this exact order. "
    "The stacking/narrowing visual shows how each layer constrains the next.")


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — What's In the Box
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "What's in the Box")

categories = [
    ("10 Personas", [
        "systems-engineer, security-auditor",
        "software-architect, devops-engineer",
        "reverse-engineer, specification-analyst",
        "workflow-arbiter, implementation-engineer",
        "test-engineer, promptkit-contributor",
    ], ACCENT),
    ("20 Protocols", [
        "Guardrails (3): anti-hallucination,\n  self-verification, operational-constraints",
        "Analysis (4): memory-safety-c, memory-safety-rust,\n  thread-safety, security-vulnerability",
        "Reasoning (13): root-cause-analysis,\n  traceability-audit, code-compliance-audit,\n  requirements-elicitation … and 9 more",
    ], ACCENT2),
    ("11 Formats", [
        "requirements-doc, design-doc,\nvalidation-plan, investigation-report",
        "implementation-plan, agent-instructions,\npipeline-spec, release-notes",
        "multi-artifact, triage-report,\npromptkit-pull-request",
    ], ACCENT3),
]

col_w = Inches(3.6)
gap = Inches(0.4)
total_w = len(categories) * col_w.inches + (len(categories) - 1) * gap.inches
sx = (13.333 - total_w) / 2

for i, (heading, items, color) in enumerate(categories):
    x = Inches(sx + i * (col_w.inches + gap.inches))
    y = Inches(1.6)
    add_rounded_rect(s, x, y, col_w, Inches(4.5), SURFACE)
    # heading
    add_text(s, x + Inches(0.2), y + Inches(0.15), col_w - Inches(0.4), Inches(0.4),
             heading, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # items — cumulative y positioning to avoid gaps with mixed line counts
    item_y = y + Inches(0.7)
    for j, item in enumerate(items):
        lines = item.count('\n') + 1
        item_h = Inches(0.25 * lines + 0.15)
        add_text(s, x + Inches(0.3), item_y,
                 col_w - Inches(0.6), item_h,
                 item, font_size=12, color=WHITE, font_name="Cascadia Code")
        item_y += Inches(0.25 * lines + 0.35)

# Bottom callout
add_rounded_rect(s, Inches(1), Inches(6.3), Inches(11.3), Inches(0.7), SURFACE)
add_text(s, Inches(1.3), Inches(6.35), Inches(10.7), Inches(0.6),
         "27 Task Templates:  investigate-bug · review-code · author-requirements-doc · "
         "author-design-doc · audit-traceability · audit-code-compliance · plan-implementation · "
         "author-pipeline · triage-issues · extend-library  … and more",
         font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)

speaker_notes(s,
    "Inventory slide. Emphasize breadth: 10 domain-expert personas, 20 protocols across "
    "3 categories (guardrails that prevent errors, analysis for domain-specific checks, "
    "reasoning for systematic approaches), 11 output format specs, and 27 task templates. "
    "All MIT-licensed, all composable.")


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Quick Start
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Quick Start  —  3 Commands")

commands = [
    ("Interactive mode", "npx @alan-jowett/promptkit",
     "Detects your LLM CLI (Copilot / Claude) and\nlaunches an interactive prompt-building session",
     ACCENT),
    ("Browse templates", "npx @alan-jowett/promptkit list",
     "Lists all available templates with descriptions\nAdd --json for machine-readable output",
     ACCENT2),
    ("Assemble a prompt", 'npx @alan-jowett/promptkit assemble investigate-bug \\\n'
     '  -p problem_description="Segfault on startup" \\\n  -o investigation.md',
     "Composes persona + protocols + format + template\ninto a single ready-to-use prompt file",
     ACCENT3),
]

for i, (label, cmd, desc, color) in enumerate(commands):
    y = Inches(1.6 + i * 1.8)
    # label
    add_text(s, Inches(0.9), y, Inches(3), Inches(0.35),
             label, font_size=18, color=color, bold=True)
    # code block
    code_rect = add_rounded_rect(s, Inches(0.9), y + Inches(0.4),
                                  Inches(7.5), Inches(0.9), CODE_BG)
    add_text(s, Inches(1.1), y + Inches(0.45), Inches(7.1), Inches(0.8),
             cmd, font_size=14, color=ACCENT2, font_name="Cascadia Code",
             line_spacing=1.4)
    # description
    add_text(s, Inches(8.8), y + Inches(0.35), Inches(4), Inches(1.0),
             desc, font_size=13, color=DIM)

add_text(s, Inches(0.9), Inches(7.0), Inches(11), Inches(0.3),
         "Prerequisites: Node.js 18+  ·  Optional: GitHub Copilot CLI or Claude Code for interactive mode",
         font_size=12, color=DIM)

speaker_notes(s,
    "Live demo opportunity. Show each command in sequence. Interactive mode is the "
    "flagship UX — it walks the user through component selection and parameter gathering. "
    "The assemble command is for automation/scripting. List is for discovery.")


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Assembly Flow (visual)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "How Assembly Works")

steps = [
    ("Persona", "# Identity\nYou are a senior\nsystems engineer …", ACCENT),
    ("Protocols", "# Reasoning\n## Anti-Hallucination\n## Root-Cause …", ACCENT2),
    ("Format", "# Output Format\nInvestigation Report\nwith sections …", ACCENT3),
    ("Template", "# Task\nInvestigate the bug:\n{{problem_description}}", ORANGE),
]

box_w = Inches(2.5)
box_h = Inches(3.2)
arrow_w = Inches(0.6)
total_w = len(steps) * box_w.inches + (len(steps) - 1) * arrow_w.inches
sx = (13.333 - total_w) / 2
y = Inches(1.8)

for i, (label, content, color) in enumerate(steps):
    x = Inches(sx + i * (box_w.inches + arrow_w.inches))
    add_rounded_rect(s, x, y, box_w, box_h, SURFACE)
    # color top bar
    s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, box_w, Inches(0.06)
                       ).fill.solid()
    s.shapes[-1].fill.fore_color.rgb = color
    s.shapes[-1].line.fill.background()
    # label
    add_text(s, x + Inches(0.15), y + Inches(0.2), box_w - Inches(0.3), Inches(0.35),
             label, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # content
    add_text(s, x + Inches(0.15), y + Inches(0.7), box_w - Inches(0.3), Inches(2.2),
             content, font_size=12, color=DIM, font_name="Cascadia Code",
             alignment=PP_ALIGN.LEFT, line_spacing=1.5)
    # arrow
    if i < len(steps) - 1:
        ax = x + box_w + Inches(0.1)
        add_text(s, ax, y + Inches(1.2), arrow_w, Inches(0.5),
                 "→", font_size=28, color=DIM, alignment=PP_ALIGN.CENTER)

# result box
result_y = Inches(5.4)
add_rounded_rect(s, Inches(2), result_y, Inches(9.3), Inches(1.2), CODE_BG)
add_text(s, Inches(2.3), result_y + Inches(0.1), Inches(8.7), Inches(1.0),
         "Output: Single composed prompt file  ─  ready to paste into any LLM\n"
         "or use as agent instructions (.github/instructions/*.md, CLAUDE.md, .cursorrules)",
         font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

speaker_notes(s,
    "Visual walkthrough of assembly. The engine reads the template's YAML frontmatter, "
    "resolves each referenced component, and concatenates them in order. Parameters like "
    "{{problem_description}} are substituted with user-provided values. The output is a "
    "single markdown file that can be pasted into any LLM or saved as persistent agent instructions.")


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Pipeline Chaining
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Pipeline Chaining: Multi-Stage Workflows")

pipeline_steps = [
    ("author-\nrequirements-doc", "requirements-\ndocument", ACCENT),
    ("author-\ndesign-doc", "design-\ndocument", ACCENT2),
    ("author-\nvalidation-plan", "validation-\nplan", ACCENT3),
]

box_w = Inches(2.8)
box_h = Inches(2.5)
arrow_w = Inches(0.8)
total_w = len(pipeline_steps) * box_w.inches + (len(pipeline_steps) - 1) * arrow_w.inches
sx = (13.333 - total_w) / 2
y = Inches(2.0)

for i, (template, artifact, color) in enumerate(pipeline_steps):
    x = Inches(sx + i * (box_w.inches + arrow_w.inches))
    add_rounded_rect(s, x, y, box_w, box_h, SURFACE)
    s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, box_w, Inches(0.06)
                       ).fill.solid()
    s.shapes[-1].fill.fore_color.rgb = color
    s.shapes[-1].line.fill.background()
    add_text(s, x + Inches(0.1), y + Inches(0.25), box_w - Inches(0.2), Inches(0.9),
             template, font_size=16, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Cascadia Code")
    # produces label
    add_text(s, x + Inches(0.1), y + Inches(1.35), box_w - Inches(0.2), Inches(0.3),
             "produces ↓", font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(1.7), box_w - Inches(0.2), Inches(0.7),
             artifact, font_size=14, color=color, alignment=PP_ALIGN.CENTER,
             font_name="Cascadia Code")
    # arrow between boxes
    if i < len(pipeline_steps) - 1:
        ax = x + box_w + Inches(0.05)
        add_text(s, ax, y + Inches(0.8), arrow_w, Inches(0.5),
                 "→", font_size=36, color=DIM, alignment=PP_ALIGN.CENTER)

# explanation
add_text(s, Inches(1.5), Inches(5.0), Inches(10), Inches(1.5),
         "Each template declares input/output contracts.\n"
         "One template's artifact becomes the next template's input.\n"
         "Build complete document lifecycles without copy-paste.",
         font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER, line_spacing=1.6)

speaker_notes(s,
    "Pipeline chaining is the power feature. Templates declare input_contract and "
    "output_contract with artifact types. The requirements doc feeds into design, "
    "design feeds into validation planning. The bootstrap engine can suggest the "
    "next stage automatically. This replaces ad-hoc copy-paste between tools.")


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Get Involved
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Get Started Today")

actions = [
    ("Try it now",
     "npx @alan-jowett/promptkit",
     "Zero install — just run it. Walks you through\ntemplate selection and prompt assembly.", ACCENT),
    ("Browse the repo",
     "github.com/microsoft/promptkit",
     "Read the templates, protocols, personas.\nUnderstand the composition model.", ACCENT2),
    ("Contribute",
     "Use the extend-library template",
     "Add your own personas, protocols, or templates.\nPromptKit eats its own dog food — the contribution\nworkflow is itself a PromptKit template.", ACCENT3),
    ("Share feedback",
     "File issues or start discussions",
     "What templates are missing? What workflows\ndo you wish PromptKit supported?", ORANGE),
]

for i, (title, cmd, desc, color) in enumerate(actions):
    y = Inches(1.55 + i * 1.35)
    add_accent_bar(s, Inches(0.9), y, height=Inches(0.9), color=color)
    add_text(s, Inches(1.2), y - Inches(0.02), Inches(4), Inches(0.35),
             title, font_size=20, color=color, bold=True)
    add_text(s, Inches(1.2), y + Inches(0.35), Inches(5), Inches(0.35),
             cmd, font_size=14, color=WHITE, font_name="Cascadia Code")
    add_text(s, Inches(7), y + Inches(0.05), Inches(5.5), Inches(0.8),
             desc, font_size=13, color=DIM)

speaker_notes(s,
    "Call to action. The barrier to entry is one npx command. Encourage the team to "
    "try it on a real task — bug investigation or code review works great as a first use. "
    "The extend-library template is meta: you use PromptKit to contribute to PromptKit.")


# ════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════
out = "docs/PromptKit-Intro.pptx"
prs.save(out)
print(f"✅ Created {out}")
