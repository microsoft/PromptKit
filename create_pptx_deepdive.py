"""Generate a PromptKit deep-dive presentation (15-18 slides, engineering audience)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import re

# ── Theme colours ──────────────────────────────────────────────
BG        = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT    = RGBColor(0x89, 0xB4, 0xFA)   # blue
ACCENT2   = RGBColor(0xA6, 0xE3, 0xA1)   # green
ACCENT3   = RGBColor(0xF9, 0xE2, 0xAF)   # yellow
ORANGE    = RGBColor(0xFA, 0xB3, 0x87)
PINK      = RGBColor(0xF3, 0x8B, 0xA8)
MAUVE     = RGBColor(0xCB, 0xA6, 0xF7)
WHITE     = RGBColor(0xCD, 0xD6, 0xF4)
DIM       = RGBColor(0x6C, 0x70, 0x86)
SURFACE   = RGBColor(0x31, 0x32, 0x44)
CODE_BG   = RGBColor(0x18, 0x18, 0x25)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── helpers ────────────────────────────────────────────────────
def dark_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_text(slide, left, top, width, height, text, *,
             font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Segoe UI", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(font_size * 0.3)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf

def add_para(tf, text, *, font_size=18, color=WHITE, bold=False,
             font_name="Segoe UI", alignment=PP_ALIGN.LEFT, space_after=None,
             line_spacing=None):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_after is not None:
        p.space_after = Pt(space_after)
    if line_spacing is not None:
        p.line_spacing = Pt(line_spacing)
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color=SURFACE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def title_strip(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)
    add_text(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.65),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.45),
                 subtitle, font_size=16, color=DIM)

def speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def code_block(slide, left, top, width, height, text, *, font_size=12,
               line_spacing=1.5):
    add_rounded_rect(slide, left, top, width, height, CODE_BG)
    add_text(slide, left + Inches(0.2), top + Inches(0.1),
             width - Inches(0.4), height - Inches(0.2),
             text, font_size=font_size, color=ACCENT2,
             font_name="Cascadia Code", line_spacing=line_spacing)

def slide_number(slide, num, total):
    add_text(slide, Inches(12.2), Inches(7.1), Inches(1), Inches(0.3),
             f"{num} / {total}", font_size=10, color=DIM,
             alignment=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 17

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
add_rect(s, Inches(0), Inches(0), W, Inches(0.08), ACCENT)

logo_size = Inches(2.5)
s.shapes.add_picture("PromptKit-logo.png", Inches(0.9), Inches(1.2),
                     logo_size, logo_size)

add_text(s, Inches(3.7), Inches(1.6), Inches(9), Inches(1.2),
         "PromptKit", font_size=64, color=WHITE, bold=True)
add_text(s, Inches(3.7), Inches(2.8), Inches(9), Inches(0.8),
         "Composable, version-controlled prompts\nfor AI-assisted software engineering",
         font_size=24, color=ACCENT)
add_text(s, Inches(3.7), Inches(4.1), Inches(9), Inches(0.5),
         "github.com/microsoft/promptkit  ·  MIT License  ·  v0.2.0",
         font_size=14, color=DIM)
add_text(s, Inches(0.9), Inches(6.2), Inches(11), Inches(0.4),
         "Deep Dive  ·  Architecture & Internals",
         font_size=14, color=DIM)
speaker_notes(s, "Title slide. This is the deep-dive version covering architecture, "
    "component internals, assembly engine, testing methodology, and pipeline chaining.")
slide_number(s, 1, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "The Problem with Ad-Hoc Prompts")

problems = [
    ("Written once, pasted everywhere", "No iteration, no improvement loop, no review process", ORANGE),
    ("No version control", "Can't diff, blame, or revert prompt changes", ACCENT3),
    ("No reuse", "Every engineer reinvents persona, guardrails, output format", ACCENT),
    ("No testing", "No way to verify quality or catch regressions", ACCENT2),
    ("Inconsistent outputs", "Same task → wildly different results across the team", PINK),
]
for i, (title, desc, color) in enumerate(problems):
    y = Inches(1.65 + i * 1.0)
    add_rect(s, Inches(0.9), y, Inches(0.06), Inches(0.6), color)
    add_text(s, Inches(1.15), y - Inches(0.03), Inches(4), Inches(0.35),
             title, font_size=20, color=WHITE, bold=True)
    add_text(s, Inches(1.15), y + Inches(0.32), Inches(10), Inches(0.35),
             desc, font_size=15, color=DIM)
speaker_notes(s, "Walk through each pain point. The core insight: prompts are the "
    "primary interface to AI tools, but we treat them as throwaway text.")
slide_number(s, 2, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — The Fix
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "The Fix: Treat Prompts as Code")

pillars = [
    ("Modular", "Compose from reusable\nbuilding blocks", ACCENT),
    ("Version-\nControlled", "Track changes in Git\nlike any code", ACCENT2),
    ("Testable", "Reference-compare\nagainst known-good", ACCENT3),
    ("Extensible", "Add your own\ncomponents via PR", ORANGE),
]
card_w = Inches(2.6)
gap = Inches(0.35)
total = len(pillars) * card_w.inches + (len(pillars) - 1) * gap.inches
sx = (13.333 - total) / 2
for i, (title, desc, color) in enumerate(pillars):
    x = Inches(sx + i * (card_w.inches + gap.inches))
    y = Inches(2.0)
    add_rounded_rect(s, x, y, card_w, Inches(2.5), SURFACE)
    add_rect(s, x, y, card_w, Inches(0.06), color)
    add_text(s, x + Inches(0.2), y + Inches(0.3), card_w - Inches(0.4), Inches(0.5),
             title, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(1.0), card_w - Inches(0.4), Inches(1.2),
             desc, font_size=15, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text(s, Inches(1.5), Inches(5.2), Inches(10), Inches(0.8),
         "PromptKit applies the same engineering discipline you use for software\n"
         "to the prompts that build it.",
         font_size=16, color=DIM, alignment=PP_ALIGN.CENTER)
speaker_notes(s, "Four engineering principles. Modularity = compose layers. "
    "Version control = full Git history. Testable = reference comparison methodology. "
    "Extensible = anyone can contribute new components.")
slide_number(s, 3, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture: 5-Layer Stack
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Architecture: The 5-Layer Composition Stack")

layers = [
    ("① Persona",    "Who the LLM becomes",              "systems-engineer, security-auditor …",  ACCENT),
    ("② Protocols",  "How it reasons + guardrails",       "anti-hallucination, root-cause-analysis …", ACCENT2),
    ("③ Format",     "Output structure & rules",          "investigation-report, design-doc …",   ACCENT3),
    ("④ Taxonomy",   "Domain classification (opt.)",      "severity tiers, risk categories",          ORANGE),
    ("⑤ Template",   "The task with {{params}}",          "investigate-bug, review-code …",         PINK),
]
for i, (name, role, examples, color) in enumerate(layers):
    y = Inches(1.55 + i * 1.05)
    bar_w = Inches(11.5 - i * 0.6)
    x = Inches((13.333 - bar_w.inches) / 2)
    add_rounded_rect(s, x, y, bar_w, Inches(0.85), SURFACE)
    add_rect(s, x, y, Inches(0.07), Inches(0.85), color)
    add_text(s, x + Inches(0.25), y + Inches(0.05), Inches(2.2), Inches(0.4),
             name, font_size=18, color=color, bold=True)
    add_text(s, x + Inches(2.5), y + Inches(0.05), Inches(3), Inches(0.4),
             role, font_size=15, color=WHITE)
    add_text(s, x + Inches(2.5), y + Inches(0.42), bar_w - Inches(2.7), Inches(0.4),
             examples, font_size=12, color=DIM)

add_text(s, Inches(1), Inches(6.9), Inches(11), Inches(0.4),
         "Templates declare which components to compose via YAML frontmatter  →  "
         "the bootstrap engine snaps them together",
         font_size=13, color=DIM, alignment=PP_ALIGN.CENTER)
speaker_notes(s, "The composition stack. Each layer constrains the next. Templates "
    "declare their dependencies in YAML frontmatter — persona, protocols, format. "
    "The narrowing visual shows how the stack funnels into a specific task.")
slide_number(s, 4, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Deep Dive: Personas
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Deep Dive: Personas", "Layer ① — Who the LLM becomes")

# Left: YAML frontmatter example
code_block(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.4),
    "---\n"
    "name: systems-engineer\n"
    "description: >\n"
    "  Senior systems engineer with deep\n"
    "  expertise in memory management,\n"
    "  concurrency, and performance.\n"
    "domain:\n"
    "  - systems programming\n"
    "  - debugging\n"
    "  - performance analysis\n"
    "tone: precise, technical, methodical\n"
    "---", font_size=12, line_spacing=1.35)

# Right: behavioral constraints
add_text(s, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.35),
         "Behavioral Constraints", font_size=20, color=ACCENT, bold=True)
constraints = [
    "Reason from first principles — trace causality, never guess",
    "Distinguish epistemic states — label KNOWN / INFERRED / ASSUMED",
    "Prefer correctness over cleverness",
    "Explicit uncertainty — say what's missing",
    "No hallucination — refuse to fabricate",
]
for i, c in enumerate(constraints):
    add_text(s, Inches(6.8), Inches(2.1 + i * 0.55), Inches(5.8), Inches(0.5),
             f"  •  {c}", font_size=13, color=WHITE)

# Bottom: all personas
add_rounded_rect(s, Inches(0.8), Inches(4.9), Inches(11.7), Inches(2.3), SURFACE)
add_text(s, Inches(1.1), Inches(5.0), Inches(11), Inches(0.35),
         "Available Personas", font_size=16, color=ACCENT3, bold=True)
personas = [
    ("systems-engineer", "Memory, concurrency, debugging"),
    ("security-auditor", "Vulnerability discovery, threat modeling"),
    ("software-architect", "System design, API contracts, tradeoffs"),
    ("devops-engineer", "CI/CD, release engineering, IaC"),
    ("reverse-engineer", "Spec extraction from existing code"),
    ("specification-analyst", "Specification drift, traceability"),
    ("workflow-arbiter", "Multi-agent workflow orchestration"),
    ("implementation-engineer", "Implementation from specs"),
    ("test-engineer", "Test authoring from specs"),
    ("promptkit-contributor", "Library architecture, conventions"),
]
for i, (name, desc) in enumerate(personas):
    col = i % 3
    row = i // 3
    x = Inches(1.1 + col * 3.9)
    y = Inches(5.45 + row * 0.42)
    add_text(s, x, y, Inches(1.8), Inches(0.35),
             name, font_size=11, color=ACCENT, bold=True, font_name="Cascadia Code")
    add_text(s, x + Inches(1.85), y, Inches(2), Inches(0.35),
             desc, font_size=11, color=DIM)

speaker_notes(s, "Personas define identity, domain expertise, tone, and behavioral "
    "constraints. The YAML frontmatter is machine-readable; the markdown body contains "
    "the actual prompt content. Key: behavioral constraints reinforce the protocols.")
slide_number(s, 5, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Deep Dive: Protocols (3 categories)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Deep Dive: Protocols", "Layer ② — How the LLM reasons")

cats = [
    ("Guardrails (3)", "Cross-cutting, apply to all tasks", [
        ("anti-hallucination", "Epistemic labeling, refusal\nto fabricate, source attribution"),
        ("self-verification", "LLM verifies its own output\nbefore finalizing"),
        ("operational-constraints", "Scoping, tool usage,\nreproducibility"),
    ], ACCENT),
    ("Analysis (4)", "Domain/language-specific checks", [
        ("memory-safety-c", "Allocation, pointers, buffers"),
        ("memory-safety-rust", "unsafe blocks, FFI, interior mut"),
        ("thread-safety", "Races, deadlocks, atomics"),
        ("security-vulnerability", "Trust boundaries, auth, crypto"),
    ], ACCENT2),
    ("Reasoning (13)", "Systematic approaches", [
        ("root-cause-analysis", "Hypothesis → evidence →\nelimination → root cause"),
        ("traceability-audit", "Cross-document requirement\ntraceability (D1–D7)"),
        ("code-compliance-audit", "Code vs specification\ncompliance (D8–D10)"),
    ], ACCENT3),
]

col_w = Inches(3.7)
gap_x = Inches(0.35)
total_w = len(cats) * col_w.inches + (len(cats) - 1) * gap_x.inches
sx = (13.333 - total_w) / 2

for ci, (cat_name, cat_desc, protocols, color) in enumerate(cats):
    x = Inches(sx + ci * (col_w.inches + gap_x.inches))
    y = Inches(1.5)
    card_h = Inches(5.5)
    add_rounded_rect(s, x, y, col_w, card_h, SURFACE)
    add_rect(s, x, y, col_w, Inches(0.06), color)
    add_text(s, x + Inches(0.15), y + Inches(0.15), col_w - Inches(0.3), Inches(0.35),
             cat_name, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y + Inches(0.55), col_w - Inches(0.3), Inches(0.3),
             cat_desc, font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)

    for pi, (pname, pdesc) in enumerate(protocols):
        py = y + Inches(1.05 + pi * 1.05)
        add_rounded_rect(s, x + Inches(0.12), py, col_w - Inches(0.24), Inches(0.95), CODE_BG)
        add_text(s, x + Inches(0.25), py + Inches(0.05), col_w - Inches(0.5), Inches(0.3),
                 pname, font_size=11, color=color, bold=True, font_name="Cascadia Code")
        add_text(s, x + Inches(0.25), py + Inches(0.38), col_w - Inches(0.5), Inches(0.55),
                 pdesc, font_size=10, color=WHITE)

    # Add "plus N more" note for categories with more protocols than shown
    shown = len(protocols)
    m = re.search(r'\((\d+)\)', cat_name)
    if m and int(m.group(1)) > shown:
        extra = int(m.group(1)) - shown
        note_y = y + Inches(1.05 + shown * 1.05 + 0.1)
        add_text(s, x + Inches(0.12), note_y, col_w - Inches(0.24), Inches(0.3),
                 f"… plus {extra} more", font_size=11, color=DIM,
                 alignment=PP_ALIGN.CENTER)

speaker_notes(s, "Three protocol categories with different scoping rules. Guardrails are "
    "cross-cutting (apply to every task). Analysis protocols are language/domain-specific. "
    "Reasoning protocols define systematic methodologies — 13 total including audit, "
    "extraction, and orchestration protocols. Separate files, not conditional blocks.")
slide_number(s, 6, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Anti-Hallucination Protocol (deep dive)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Spotlight: Anti-Hallucination Protocol",
            "The most important guardrail — prevents fabrication")

rules = [
    ("Epistemic Labeling",
     "Every claim tagged: KNOWN (directly stated), INFERRED (reasoned with "
     "explicit chain), or ASSUMED (flagged with justification).\n"
     "If ASSUMED > ~30% → stop and request more context.",
     ACCENT),
    ("Refusal to Fabricate",
     "Do NOT invent function names, API signatures, config values, "
     "file paths, or version numbers.\n"
     'Use [UNKNOWN: <what is missing>] as placeholder.',
     ORANGE),
    ("Uncertainty Disclosure",
     "Enumerate multiple interpretations rather than choosing silently. "
     'State: "Low confidence — depends on [assumption]. Verify by [action]."',
     ACCENT3),
    ("Source Attribution",
     'Indicate provenance: "per requirements doc §3.2" or "line 42 of auth.c".\n'
     "Do NOT cite sources not provided.",
     ACCENT2),
    ("Scope Boundaries",
     "If question falls outside provided context, say so explicitly.\n"
     "List what additional information is needed.",
     PINK),
]

for i, (title, desc, color) in enumerate(rules):
    y = Inches(1.5 + i * 1.12)
    add_rect(s, Inches(0.8), y, Inches(0.06), Inches(0.85), color)
    add_text(s, Inches(1.1), y - Inches(0.02), Inches(3), Inches(0.3),
             title, font_size=18, color=color, bold=True)
    add_text(s, Inches(4.0), y + Inches(0.02), Inches(8.5), Inches(0.85),
             desc, font_size=12, color=WHITE, line_spacing=1.4)

speaker_notes(s, "Deep dive into the anti-hallucination protocol. This is the single "
    "most impactful guardrail. Epistemic labeling forces the LLM to be honest about "
    "what it knows vs infers vs assumes. The 30% threshold prevents assumption-heavy outputs.")
slide_number(s, 7, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Root Cause Analysis Protocol
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Spotlight: Root Cause Analysis Protocol",
            "5-phase structured reasoning methodology")

phases = [
    ("Phase 1", "Symptom\nCharacterization",
     "Observed vs expected, timeline,\nblast radius, deterministic?", ACCENT),
    ("Phase 2", "Hypothesis\nGeneration",
     "≥3 ranked hypotheses before\ninvestigating any. Include 1 non-obvious.", ACCENT2),
    ("Phase 3", "Evidence &\nElimination",
     "Minimal investigation per hypothesis.\nCONFIRMED / ELIMINATED / INCONCLUSIVE", ACCENT3),
    ("Phase 4", "Root Cause\nIdentification",
     "Distinguish root from proximate.\nTrace full causal chain.", ORANGE),
    ("Phase 5", "Remediation",
     "Fix root, not symptom.\nSecondary fixes + risk assessment.", PINK),
]

box_w = Inches(2.1)
gap = Inches(0.3)
total_w = len(phases) * box_w.inches + (len(phases) - 1) * gap.inches
sx = (13.333 - total_w) / 2
y = Inches(1.7)

for i, (label, title, desc, color) in enumerate(phases):
    x = Inches(sx + i * (box_w.inches + gap.inches))
    add_rounded_rect(s, x, y, box_w, Inches(3.5), SURFACE)
    add_rect(s, x, y, box_w, Inches(0.06), color)
    add_text(s, x + Inches(0.1), y + Inches(0.2), box_w - Inches(0.2), Inches(0.25),
             label, font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(0.5), box_w - Inches(0.2), Inches(0.7),
             title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(1.5), box_w - Inches(0.2), Inches(1.6),
             desc, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER, line_spacing=1.4)
    if i < len(phases) - 1:
        add_text(s, x + box_w + Inches(0.02), y + Inches(1.4), Inches(0.3), Inches(0.4),
                 "→", font_size=24, color=DIM, alignment=PP_ALIGN.CENTER)

# Key insight callout
add_rounded_rect(s, Inches(1.5), Inches(5.6), Inches(10.3), Inches(1.2), CODE_BG)
add_text(s, Inches(1.8), Inches(5.7), Inches(9.7), Inches(1.0),
         'Key rule: "If we fix only the proximate cause, will the root cause\n'
         'produce other failures?" → If yes, the fix is incomplete.\n\n'
         'Example: Proximate = "null deref at line 42" · Root = "init silently fails when config missing"',
         font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

speaker_notes(s, "The RCA protocol enforces systematic debugging. The ≥3 hypothesis rule "
    "prevents anchoring bias. The root vs proximate distinction is critical — most quick fixes "
    "only address the proximate cause and leave the root cause to produce more bugs.")
slide_number(s, 8, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — All Output Formats (table)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Deep Dive: Output Formats",
            "Layer ③ — 11 structured deliverable types")

formats = [
    ("requirements-doc", "Numbered REQ-IDs + acceptance criteria", "requirements-document"),
    ("design-doc", "Architecture, APIs, tradeoff analysis", "design-document"),
    ("validation-plan", "Test cases + traceability matrix", "validation-plan"),
    ("investigation-report", "Findings, RCA, remediation", "investigation-report"),
    ("implementation-plan", "Task breakdown + dependencies", "implementation-plan"),
    ("agent-instructions", ".github/instructions/, CLAUDE.md, .cursorrules", "agent-instructions"),
    ("pipeline-spec", "YAML pipelines + design rationale", "pipeline-spec"),
    ("release-notes", "Changelog, breaking changes, upgrade path", "release-notes"),
    ("multi-artifact", "Multiple structured files", "multi-artifact"),
    ("triage-report", "Issue/PR triage with priority ranking", "triage-report"),
    ("promptkit-pull-request", "PromptKit contribution PRs", "promptkit-pull-request"),
]

# Centered full-width table
table_x = Inches(1.5)
table_w = Inches(10.3)

# header
add_rounded_rect(s, table_x, Inches(1.7), table_w, Inches(0.45), ACCENT)
add_text(s, table_x + Inches(0.2), Inches(1.75), Inches(2.5), Inches(0.35),
         "Format", font_size=14, color=BG, bold=True)
add_text(s, table_x + Inches(2.8), Inches(1.75), Inches(4.5), Inches(0.35),
         "Description", font_size=14, color=BG, bold=True)
add_text(s, table_x + Inches(7.5), Inches(1.75), Inches(2.5), Inches(0.35),
         "Artifact Type", font_size=14, color=BG, bold=True)

for i, (name, desc, artifact) in enumerate(formats):
    y = Inches(2.25 + i * 0.42)
    bg = SURFACE if i % 2 == 0 else CODE_BG
    add_rounded_rect(s, table_x, y, table_w, Inches(0.38), bg)
    add_text(s, table_x + Inches(0.2), y + Inches(0.04), Inches(2.5), Inches(0.3),
             name, font_size=12, color=ACCENT, font_name="Cascadia Code")
    add_text(s, table_x + Inches(2.8), y + Inches(0.04), Inches(4.5), Inches(0.3),
             desc, font_size=12, color=WHITE)
    add_text(s, table_x + Inches(7.5), y + Inches(0.04), Inches(2.5), Inches(0.3),
             artifact, font_size=10, color=DIM, font_name="Cascadia Code")

# Footer note
add_text(s, Inches(1.5), Inches(7.0), Inches(10.3), Inches(0.3),
         "Each format declares a produces field → enables pipeline chaining between templates",
         font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)

speaker_notes(s, "Overview of all 11 output formats. Each format defines the structure "
    "of the deliverable — section ordering, required content, naming conventions. "
    "The 'produces' field enables pipeline chaining between templates.")
slide_number(s, 9, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Spotlight: Investigation Report Format
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Spotlight: Investigation Report",
            "The most common output format — 9 structured sections")

# Left: format frontmatter
code_block(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(1.8),
    "---\n"
    "name: investigation-report\n"
    "type: format\n"
    "produces: investigation-report\n"
    "description: >\n"
    "  Structured findings with root cause analysis\n"
    "---", font_size=12, line_spacing=1.3)

# Left: section list
add_text(s, Inches(0.8), Inches(3.6), Inches(5.5), Inches(0.35),
         "Required Sections", font_size=18, color=ACCENT, bold=True)
sections = [
    ("1. Executive Summary", "2-4 sentence overview of findings"),
    ("2. Problem Statement", "What was reported and by whom"),
    ("3. Investigation Scope", "Boundaries + limitations of the analysis"),
    ("4. Findings (F-NNN)", "Severity-ordered, classified findings"),
    ("5. Root Cause Analysis", "Full causal chain, root vs proximate"),
    ("6. Remediation Plan", "Priority table with effort estimates"),
    ("7. Prevention", "Code, process, and tooling changes"),
    ("8. Open Questions", "What remains unknown or unverified"),
    ("9. Revision History", "Audit trail of report changes"),
]
for i, (sec, desc) in enumerate(sections):
    y = Inches(4.1 + i * 0.36)
    add_text(s, Inches(0.8), y, Inches(2.8), Inches(0.35),
             sec, font_size=12, color=WHITE, bold=True)
    add_text(s, Inches(3.6), y, Inches(4), Inches(0.35),
             desc, font_size=12, color=DIM)

# Right: key rules
add_text(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.35),
         "Format Rules", font_size=18, color=ACCENT3, bold=True)
rules = [
    ("Severity ordering", "Findings presented Critical → High →\nMedium → Low → Informational"),
    ("F-NNN identifiers", "Every finding gets a unique ID\nfor cross-referencing"),
    ("Causal chains", "Root cause must trace back to\nthe proximate symptom"),
    ("Remediation table", "Each finding maps to a fix\nwith priority and effort"),
    ("Non-goals section", "Explicitly state what is NOT\nin scope for this report"),
]
for i, (title, desc) in enumerate(rules):
    y = Inches(2.0 + i * 1.0)
    add_rect(s, Inches(7.0), y, Inches(0.06), Inches(0.75), ACCENT)
    add_text(s, Inches(7.25), y, Inches(5.5), Inches(0.3),
             title, font_size=14, color=ACCENT, bold=True)
    add_text(s, Inches(7.25), y + Inches(0.3), Inches(5.5), Inches(0.5),
             desc, font_size=11, color=WHITE, line_spacing=1.3)

speaker_notes(s, "Deep dive into the investigation-report format. The 9-section structure "
    "ensures consistency across all bug investigations and security audits. Key: findings "
    "use F-NNN IDs, severity ordering, and the root cause section must trace the full "
    "causal chain — not just the proximate symptom.")
slide_number(s, 10, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Deep Dive: Templates (frontmatter)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Deep Dive: Template Anatomy", "Layer ⑤ — The task declaration")

# Left: real frontmatter
code_block(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0),
    "---\n"
    "name: investigate-bug\n"
    "description: >\n"
    "  Systematically investigate a bug.\n"
    "  Apply root cause analysis and produce\n"
    "  an investigation report.\n"
    "persona: systems-engineer\n"
    "protocols:\n"
    "  - guardrails/anti-hallucination\n"
    "  - guardrails/self-verification\n"
    "  - guardrails/operational-constraints\n"
    "  - reasoning/root-cause-analysis\n"
    "format: investigation-report\n"
    "params:\n"
    "  problem_description: \"Bug description\"\n"
    "  code_context: \"Relevant code/logs\"\n"
    "  environment: \"OS, runtime, build\"\n"
    "input_contract: null\n"
    "output_contract:\n"
    "  type: investigation-report\n"
    "---", font_size=11, line_spacing=1.3)

# Right: field explanations
fields = [
    ("persona:", "References one persona file.\nBecomes the # Identity section.", ACCENT),
    ("protocols:", "List of category/name paths.\nEach loaded as # Reasoning Protocol.", ACCENT2),
    ("format:", "References one format file.\nBecomes the # Output Format section.", ACCENT3),
    ("params:", "Key-value map → {{param}} placeholders\nin template body get substituted.", ORANGE),
    ("input_contract:", "What artifact this template consumes.\nnull = no prerequisite.", PINK),
    ("output_contract:", "What artifact this template produces.\nEnables pipeline chaining.", MAUVE),
]

for i, (field, desc, color) in enumerate(fields):
    y = Inches(1.5 + i * 0.9)
    add_text(s, Inches(6.7), y, Inches(2.5), Inches(0.3),
             field, font_size=15, color=color, bold=True, font_name="Cascadia Code")
    add_text(s, Inches(6.7), y + Inches(0.32), Inches(6), Inches(0.55),
             desc, font_size=12, color=WHITE, line_spacing=1.4)

speaker_notes(s, "Template frontmatter is the composition declaration. The key insight: "
    "protocols use category/name paths (guardrails/anti-hallucination) while the manifest "
    "uses short names (anti-hallucination). CI validates they stay in sync.")
slide_number(s, 11, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — The Manifest
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "The Manifest: Source of Truth", "manifest.yaml — the component registry")

# Left: manifest structure
code_block(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2),
    "# manifest.yaml\n"
    "version: \"0.2.0\"\n"
    "personas:\n"
    "  - name: systems-engineer\n"
    "    path: personas/systems-engineer.md\n"
    "protocols:\n"
    "  guardrails:\n"
    "    - name: anti-hallucination\n"
    "      path: protocols/guardrails/...\n"
    "  analysis:\n"
    "    - name: memory-safety-c\n"
    "  reasoning:\n"
    "    - name: root-cause-analysis\n"
    "formats:\n"
    "  - name: investigation-report\n"
    "    produces: investigation-report\n"
    "templates:\n"
    "  investigation:\n"
    "    - name: investigate-bug\n"
    "      protocols: [anti-hallucination, ...]\n"
    "      format: investigation-report\n"
    "pipelines:\n"
    "  document-lifecycle:\n"
    "    stages:\n"
    "      - template: author-requirements-doc\n"
    "        produces: requirements-document", font_size=10, line_spacing=1.3)

# Right: key rules
add_text(s, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.35),
         "Key Rules", font_size=20, color=ACCENT, bold=True)

rules = [
    ("Protocol naming duality",
     "Templates use category paths:\n  guardrails/anti-hallucination\n"
     "Manifest uses short names:\n  anti-hallucination\n"
     "CI enforces they stay in sync.", ACCENT),
    ("Path resolution",
     "Every component has a path field\npointing to its .md source file.\n"
     "Assembly engine reads these paths.", ACCENT2),
    ("Pipeline stages",
     "produces/consumes fields enable\nthe bootstrap engine to validate\n"
     "chaining and suggest next stages.", ACCENT3),
    ("CI validation",
     "tests/validate-manifest.py checks\nthat manifest protocols match\n"
     "template frontmatter protocols.", ORANGE),
]
for i, (title, desc, color) in enumerate(rules):
    y = Inches(2.1 + i * 1.25)
    add_rect(s, Inches(6.8), y, Inches(0.06), Inches(1.0), color)
    add_text(s, Inches(7.05), y, Inches(5.5), Inches(0.3),
             title, font_size=14, color=color, bold=True)
    add_text(s, Inches(7.05), y + Inches(0.3), Inches(5.5), Inches(0.7),
             desc, font_size=11, color=WHITE, font_name="Cascadia Code", line_spacing=1.3)

speaker_notes(s, "The manifest is the component registry. Critical detail: protocols have "
    "a naming duality — templates use category/name paths, manifest uses short names. "
    "CI (validate-manifest.py) enforces sync. This is the only automated check in the repo.")
slide_number(s, 12, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Bootstrap Flow
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "The Bootstrap Flow", "bootstrap.md — the meta-prompt that drives composition")

steps = [
    ("1", "Understand\nthe task", "Ask user what they\nwant to accomplish", ACCENT),
    ("2", "Read\nmanifest", "Discover all available\ncomponents", ACCENT2),
    ("3", "Select\ntemplate", "Match task to\nbest template", ACCENT3),
    ("4", "Check\nmode", "Interactive → execute\nSingle-shot → file", ORANGE),
    ("5", "Collect\nparams", "Gather template-\nspecific inputs", PINK),
    ("6", "Assemble", "Compose layers\nin order", MAUVE),
    ("7", "Output", "Write prompt file\nor execute directly", ACCENT),
]

box_w = Inches(1.55)
box_h = Inches(2.5)
gap = Inches(0.15)
total_w = len(steps) * box_w.inches + (len(steps) - 1) * gap.inches
sx = (13.333 - total_w) / 2

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(sx + i * (box_w.inches + gap.inches))
    y = Inches(1.7)
    add_rounded_rect(s, x, y, box_w, box_h, SURFACE)
    add_rect(s, x, y, box_w, Inches(0.05), color)
    # number circle
    add_text(s, x + Inches(0.05), y + Inches(0.15), box_w - Inches(0.1), Inches(0.3),
             num, font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.05), y + Inches(0.55), box_w - Inches(0.1), Inches(0.7),
             title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.05), y + Inches(1.4), box_w - Inches(0.1), Inches(1.0),
             desc, font_size=10, color=DIM, alignment=PP_ALIGN.CENTER, line_spacing=1.4)
    if i < len(steps) - 1:
        add_text(s, x + box_w - Inches(0.05), y + Inches(1.0), Inches(0.35), Inches(0.3),
                 "›", font_size=20, color=DIM, alignment=PP_ALIGN.CENTER)

# Two modes callout
add_rounded_rect(s, Inches(0.8), Inches(4.6), Inches(5.7), Inches(2.3), CODE_BG)
add_text(s, Inches(1.0), Inches(4.7), Inches(5.3), Inches(0.3),
         "Single-Shot Mode (default)", font_size=15, color=ACCENT, bold=True)
add_text(s, Inches(1.0), Inches(5.1), Inches(5.3), Inches(1.0),
         "Gather params → assemble prompt → write .md file\n"
         "User pastes into any LLM session\n\n"
         "Output options:\n"
         "  (a) Raw prompt file\n"
         "  (b) Agent instructions (.github/instructions/,\n"
         "       CLAUDE.md, .cursorrules)",
         font_size=11, color=WHITE, font_name="Cascadia Code", line_spacing=1.3)

add_rounded_rect(s, Inches(6.8), Inches(4.6), Inches(5.7), Inches(2.3), CODE_BG)
add_text(s, Inches(7.0), Inches(4.7), Inches(5.3), Inches(0.3),
         "Interactive Mode", font_size=15, color=ACCENT2, bold=True)
add_text(s, Inches(7.0), Inches(5.1), Inches(5.3), Inches(1.0),
         "Load components → execute directly in session\n"
         "Reasoning-and-challenge phase before output\n\n"
         "Template declares: mode: interactive\n"
         "Used by: extend-library, interactive-design\n"
         "No file written — LLM acts on the prompt live",
         font_size=11, color=WHITE, font_name="Cascadia Code", line_spacing=1.3)

speaker_notes(s, "The bootstrap flow is the core UX. bootstrap.md is itself a meta-prompt — "
    "it instructs the LLM how to use PromptKit. Two modes: single-shot writes a file you "
    "paste elsewhere; interactive executes the prompt live in the current session.")
slide_number(s, 13, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — Assembly Engine Internals
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Assembly Engine Internals", "cli/lib/assemble.js — how the sausage is made")

# Assembly order diagram
order_items = [
    ("# Identity", "Load persona .md, strip frontmatter", ACCENT),
    ("# Reasoning Protocols", "Load each protocol .md, join with ---", ACCENT2),
    ("# Output Format", "Load format .md, strip frontmatter", ACCENT3),
    ("# Task", "Load template .md, substitute {{params}}", ORANGE),
]

for i, (section, desc, color) in enumerate(order_items):
    y = Inches(1.6 + i * 1.15)
    add_rounded_rect(s, Inches(0.8), y, Inches(5.5), Inches(0.95), SURFACE)
    add_rect(s, Inches(0.8), y, Inches(0.07), Inches(0.95), color)
    add_text(s, Inches(1.1), y + Inches(0.08), Inches(2.5), Inches(0.35),
             section, font_size=16, color=color, bold=True, font_name="Cascadia Code")
    add_text(s, Inches(1.1), y + Inches(0.48), Inches(4.8), Inches(0.4),
             desc, font_size=13, color=WHITE)
    if i < len(order_items) - 1:
        add_text(s, Inches(3.3), y + Inches(0.95), Inches(1), Inches(0.3),
                 "↓", font_size=18, color=DIM, alignment=PP_ALIGN.CENTER)

# Right: key functions
add_text(s, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.3),
         "Key Processing Steps", font_size=18, color=ACCENT, bold=True)

funcs = [
    ("stripFrontmatter()",
     "Removes YAML --- block from\neach component before assembly", ACCENT),
    ("Strip HTML comments",
     "Removes SPDX license headers\n(<!-- ... --> blocks)", ACCENT2),
    ("loadComponent()",
     "Reads file, strips frontmatter +\ncomments, returns clean body", ACCENT3),
    ("Parameter substitution",
     "Simple {{key}} → value replacement\nvia split/join (no recursion)", ORANGE),
    ("Section joining",
     "Sections separated by markdown\nhorizontal rule (---)", PINK),
]

for i, (name, desc, color) in enumerate(funcs):
    y = Inches(2.1 + i * 0.95)
    add_text(s, Inches(6.8), y, Inches(3), Inches(0.25),
             name, font_size=13, color=color, bold=True, font_name="Cascadia Code")
    add_text(s, Inches(6.8), y + Inches(0.28), Inches(5.5), Inches(0.55),
             desc, font_size=11, color=WHITE, line_spacing=1.3)

# Bottom: output
add_rounded_rect(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8), CODE_BG)
add_text(s, Inches(1.1), Inches(6.3), Inches(11.1), Inches(0.6),
         "Output: Single markdown string  →  # Identity\\n---\\n# Reasoning Protocols\\n---\\n"
         "# Output Format\\n---\\n# Task",
         font_size=12, color=ACCENT, alignment=PP_ALIGN.CENTER, font_name="Cascadia Code")

speaker_notes(s, "The assembly engine is deliberately simple. It reads the manifest to "
    "resolve component paths, loads each .md file, strips YAML frontmatter and SPDX headers, "
    "wraps in section headers, joins with horizontal rules, and substitutes parameters. "
    "No template language, no recursion, no magic — just concatenation in a specific order.")
slide_number(s, 14, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 15 — Pipeline Chaining
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Pipeline Chaining: Multi-Stage Workflows",
            "Templates connected by input/output contracts")

pipeline_steps = [
    ("author-\nrequirements-doc", "requirements-\ndocument", ACCENT),
    ("author-\ndesign-doc", "design-\ndocument", ACCENT2),
    ("author-\nvalidation-plan", "validation-\nplan", ACCENT3),
]

box_w = Inches(2.8)
box_h = Inches(2.5)
arrow_w = Inches(0.8)
total_w = len(pipeline_steps) * box_w.inches + (len(pipeline_steps) - 1) * arrow_w.inches
sx = (13.333 - total_w) / 2
y = Inches(1.7)

for i, (template, artifact, color) in enumerate(pipeline_steps):
    x = Inches(sx + i * (box_w.inches + arrow_w.inches))
    add_rounded_rect(s, x, y, box_w, box_h, SURFACE)
    add_rect(s, x, y, box_w, Inches(0.06), color)
    add_text(s, x + Inches(0.1), y + Inches(0.25), box_w - Inches(0.2), Inches(0.9),
             template, font_size=16, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Cascadia Code")
    add_text(s, x + Inches(0.1), y + Inches(1.35), box_w - Inches(0.2), Inches(0.3),
             "produces ↓", font_size=12, color=DIM, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(1.7), box_w - Inches(0.2), Inches(0.7),
             artifact, font_size=14, color=color, alignment=PP_ALIGN.CENTER,
             font_name="Cascadia Code")
    if i < len(pipeline_steps) - 1:
        ax = x + box_w + Inches(0.05)
        add_text(s, ax, y + Inches(0.8), arrow_w, Inches(0.5),
                 "→", font_size=36, color=DIM, alignment=PP_ALIGN.CENTER)

# Contract mechanism
add_rounded_rect(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.4), CODE_BG)
add_text(s, Inches(1.1), Inches(4.7), Inches(5), Inches(0.3),
         "How contracts work:", font_size=16, color=ACCENT, bold=True)

contract_lines = [
    ("1.", "Template declares", "output_contract: { type: requirements-document }", ACCENT),
    ("2.", "Next template declares", "input_contract: { type: requirements-document }", ACCENT2),
    ("3.", "Manifest defines pipeline", "stages: [{produces: ...}, {consumes: ..., produces: ...}]", ACCENT3),
    ("4.", "Bootstrap validates", "Previous output matches next input before offering template", ORANGE),
    ("5.", "Artifacts flow forward", "Output of stage N becomes input parameter of stage N+1", PINK),
]

for i, (num, label, detail, color) in enumerate(contract_lines):
    y = Inches(5.15 + i * 0.35)
    add_text(s, Inches(1.1), y, Inches(0.4), Inches(0.3),
             num, font_size=12, color=color, bold=True)
    add_text(s, Inches(1.5), y, Inches(2.2), Inches(0.3),
             label, font_size=12, color=WHITE)
    add_text(s, Inches(3.8), y, Inches(8.3), Inches(0.3),
             detail, font_size=11, color=DIM, font_name="Cascadia Code")

speaker_notes(s, "Pipeline chaining is the power feature. Templates declare input_contract "
    "and output_contract with artifact types. The manifest defines pipeline stages. "
    "The bootstrap engine validates contracts before offering the next template. "
    "This replaces ad-hoc copy-paste between tools with a structured workflow.")
slide_number(s, 15, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 16 — Testing Methodology
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Testing: Reference Comparison Methodology",
            "How do you test prompts? Compare against known-good references.")

# Flow
flow_steps = [
    ("Hand-craft\nreference prompt", ACCENT),
    ("Generate via\nPromptKit", ACCENT2),
    ("Structured\ngap analysis", ACCENT3),
    ("Feed gaps back\nas improvements", ORANGE),
]
box_w = Inches(2.5)
gap = Inches(0.5)
total_w = len(flow_steps) * box_w.inches + (len(flow_steps) - 1) * gap.inches
sx = (13.333 - total_w) / 2

for i, (label, color) in enumerate(flow_steps):
    x = Inches(sx + i * (box_w.inches + gap.inches))
    add_rounded_rect(s, x, Inches(1.6), box_w, Inches(1.1), SURFACE)
    add_rect(s, x, Inches(1.6), box_w, Inches(0.05), color)
    add_text(s, x + Inches(0.1), Inches(1.75), box_w - Inches(0.2), Inches(0.8),
             label, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(flow_steps) - 1:
        add_text(s, x + box_w + Inches(0.08), Inches(1.9), Inches(0.4), Inches(0.4),
                 "→", font_size=24, color=DIM, alignment=PP_ALIGN.CENTER)

# 5 dimensions
add_text(s, Inches(0.8), Inches(3.1), Inches(11), Inches(0.35),
         "5 Gap Analysis Dimensions", font_size=20, color=ACCENT, bold=True)

dimensions = [
    ("Task Framing", "Goal statement, success criteria,\nnon-goals, context definition",
     "Is the objective clearly stated?\nIs scope explicitly bounded?", ACCENT),
    ("Reasoning\nMethodology", "Systematic analysis, hypothesis\ngeneration, evidence requirements",
     "Is anti-hallucination present?\nAre multiple hypotheses required?", ACCENT2),
    ("Output\nSpecification", "Format structure, artifacts,\ntaxonomy, severity ranking",
     "Is the deliverable structure defined?\nAre classification schemes provided?", ACCENT3),
    ("Operational\nGuidance", "Scoping strategy, tool usage,\nstep-by-step plan",
     "Does it guide how to scope work?\nIs a procedural plan provided?", ORANGE),
    ("Quality\nAssurance", "Self-verification, sampling,\ncoverage, consistency",
     "Must LLM verify its output?\nMust it document what it examined?", PINK),
]

col_w = Inches(2.25)
gap_x = Inches(0.2)
total_w = len(dimensions) * col_w.inches + (len(dimensions) - 1) * gap_x.inches
sx = (13.333 - total_w) / 2

for i, (title, checks, questions, color) in enumerate(dimensions):
    x = Inches(sx + i * (col_w.inches + gap_x.inches))
    add_rounded_rect(s, x, Inches(3.6), col_w, Inches(3.5), SURFACE)
    add_rect(s, x, Inches(3.6), col_w, Inches(0.05), color)
    add_text(s, x + Inches(0.1), Inches(3.75), col_w - Inches(0.2), Inches(0.55),
             title, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), Inches(4.4), col_w - Inches(0.2), Inches(1.0),
             checks, font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER, line_spacing=1.3)
    add_text(s, x + Inches(0.1), Inches(5.5), col_w - Inches(0.2), Inches(1.2),
             questions, font_size=9, color=DIM, alignment=PP_ALIGN.CENTER, line_spacing=1.3)

speaker_notes(s, "Testing prompts is hard because output is non-deterministic. The reference "
    "comparison methodology sidesteps this: compare prompt structure, not LLM output. "
    "Each dimension gets ✅ Covered / ⚠️ Partial / ❌ Missing classification. "
    "Gaps feed back into the library as improvements.")
slide_number(s, 16, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 17 — Get Involved
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s)
title_strip(s, "Get Started Today")

actions = [
    ("Try it now",
     "npx @alan-jowett/promptkit",
     "Zero install — walks you through template\nselection and prompt assembly.", ACCENT),
    ("Browse the repo",
     "github.com/microsoft/promptkit",
     "Read the templates, protocols, personas.\nUnderstand the composition model.", ACCENT2),
    ("Contribute",
     "Use the extend-library template",
     "PromptKit eats its own dog food — the\ncontribution workflow is a PromptKit template.", ACCENT3),
    ("Share feedback",
     "File issues or start discussions",
     "What templates are missing? What workflows\ndo you wish PromptKit supported?", ORANGE),
]

for i, (title, cmd, desc, color) in enumerate(actions):
    y = Inches(1.55 + i * 1.35)
    add_rect(s, Inches(0.9), y, Inches(0.06), Inches(0.9), color)
    add_text(s, Inches(1.2), y - Inches(0.02), Inches(4), Inches(0.35),
             title, font_size=20, color=color, bold=True)
    add_text(s, Inches(1.2), y + Inches(0.35), Inches(5), Inches(0.35),
             cmd, font_size=14, color=WHITE, font_name="Cascadia Code")
    add_text(s, Inches(7), y + Inches(0.05), Inches(5.5), Inches(0.8),
             desc, font_size=13, color=DIM)

speaker_notes(s, "Call to action. The barrier to entry is one npx command. "
    "The extend-library template is meta: you use PromptKit to contribute to PromptKit.")
slide_number(s, 17, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
out = "docs/PromptKit-DeepDive.pptx"
prs.save(out)
print(f"✅ Created {out} ({TOTAL_SLIDES} slides)")
