#!/usr/bin/env python3
"""Generate PromptKit Office Hour presentation using python-pptx."""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === Theme: Corporate Blue ===
BG_COLOR = RGBColor(0xF5, 0xF7, 0xFA)
TEXT_COLOR = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT_COLOR = RGBColor(0x00, 0x52, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_ACCENT = RGBColor(0xE8, 0xEF, 0xF7)
DARK_ACCENT = RGBColor(0x00, 0x3D, 0x99)
SUBTLE_TEXT = RGBColor(0x66, 0x66, 0x66)
DIVIDER_BG = RGBColor(0x00, 0x52, 0xCC)
RED_ACCENT = RGBColor(0xCC, 0x00, 0x00)
GREEN_ACCENT = RGBColor(0x00, 0x7A, 0x33)

FONT_HEADING = "Segoe UI Semibold"
FONT_BODY = "Segoe UI"
FONT_MONO = "Cascadia Code"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

OUTPUT_DIR = r"F:\promptkit_slides_office_hour"
LOGO_PATH = r"F:\promptkit\PromptKit-logo.png"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
                font_size=18, color=TEXT_COLOR, bold=False, alignment=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, bullets, left=Inches(0.8), top=Inches(1.8),
                     width=Inches(11.7), height=Inches(5.0), font_size=18,
                     color=TEXT_COLOR, spacing=Pt(4), group_spacing=Pt(14)):
    """Add bulleted text to a slide.
    
    Spacing model: empty strings in bullets become group gaps via
    group_spacing on the *preceding* paragraph, rather than creating
    empty paragraphs that consume full line height.
    """
    # Filter out empty strings; instead mark prior item for extra spacing
    processed = []  # list of (text, extra_gap_after)
    for b in bullets:
        if b.strip() == "":
            if processed:
                processed[-1] = (processed[-1][0], True)
        else:
            processed.append((b, False))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, (bullet, gap_after) in enumerate(processed):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle bold prefix (text before first non-bold segment)
        if "**" in bullet:
            parts = bullet.split("**")
            for j, part in enumerate(parts):
                if not part:
                    continue
                run = p.add_run()
                run.text = part
                run.font.name = FONT_BODY
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.bold = (j % 2 == 1)
        else:
            run = p.add_run()
            run.text = bullet
            run.font.name = FONT_BODY
            run.font.size = Pt(font_size)
            run.font.color.rgb = color

        p.space_after = group_spacing if gap_after else spacing
        p.level = 0

    # No spacing after the last paragraph
    if processed:
        tf.paragraphs[-1].space_after = Pt(0)

    return txBox


def add_shape_rect(slide, left, top, width, height, fill_color, text="",
                   font_size=14, font_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = FONT_BODY
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
    return shape


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def make_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_COLOR)

    # Accent bar at top
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), ACCENT_COLOR)

    # Logo
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(4.667), Inches(0.6), Inches(4), Inches(2.5))

    # Title
    add_textbox(slide, Inches(1), Inches(3.3), Inches(11.333), Inches(1.2),
                "Introduction to Structured Prompts\nwith PromptKit",
                FONT_HEADING, 40, ACCENT_COLOR, True, PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(1), Inches(4.6), Inches(11.333), Inches(0.6),
                "Bringing software engineering discipline to AI prompts",
                FONT_BODY, 22, SUBTLE_TEXT, False, PP_ALIGN.CENTER)

    # Speaker
    add_textbox(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.5),
                "Alan Jowett",
                FONT_BODY, 20, TEXT_COLOR, False, PP_ALIGN.CENTER)

    # Bottom bar
    add_shape_rect(slide, Inches(0), Inches(7.2), SLIDE_WIDTH, Inches(0.3), ACCENT_COLOR)

    set_notes(slide, "Welcome everyone. I'm Alan Jowett, and today I want to talk about something I think is going to change how we work with AI tools — structured, engineered prompts.")
    return slide


def make_content_slide(prs, title, bullets, notes, highlight_first=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_COLOR)

    # Top accent bar
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_COLOR)

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0),
                title, FONT_HEADING, 32, ACCENT_COLOR, True)

    # Underline accent
    add_shape_rect(slide, Inches(0.8), Inches(1.35), Inches(3), Inches(0.06), ACCENT_COLOR)

    # Auto-size: count real lines (non-empty) to pick font size
    real_lines = sum(1 for b in bullets if b.strip() != "")
    if real_lines > 14:
        fs = 14
    elif real_lines > 11:
        fs = 16
    else:
        fs = 18

    # Bullets
    add_bullet_slide(slide, bullets, font_size=fs)

    set_notes(slide, notes)
    return slide


def make_section_divider(prs, title, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DIVIDER_BG)

    # Center title
    add_textbox(slide, Inches(1), Inches(2.5), Inches(11.333), Inches(2.5),
                title, FONT_HEADING, 44, WHITE, True, PP_ALIGN.CENTER,
                MSO_ANCHOR.MIDDLE)

    # Decorative lines
    add_shape_rect(slide, Inches(5), Inches(2.2), Inches(3.333), Inches(0.06), WHITE)
    add_shape_rect(slide, Inches(5), Inches(5.3), Inches(3.333), Inches(0.06), WHITE)

    set_notes(slide, notes)
    return slide


def make_diagram_slide(prs, title, diagram_elements, notes, description=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_COLOR)

    # Top accent bar
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_COLOR)

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0),
                title, FONT_HEADING, 32, ACCENT_COLOR, True)

    add_shape_rect(slide, Inches(0.8), Inches(1.35), Inches(3), Inches(0.06), ACCENT_COLOR)

    # Build diagram from elements (list of dicts with label, sublabel, y_offset)
    box_width = Inches(9)
    box_height = Inches(0.85)
    left = Inches(2.167)

    colors = [ACCENT_COLOR, RGBColor(0x00, 0x6A, 0xCC), RGBColor(0x00, 0x82, 0xCC),
              RGBColor(0x00, 0x9A, 0xCC), DARK_ACCENT]

    for i, elem in enumerate(diagram_elements):
        y = Inches(1.7) + Inches(i * 1.05)
        color = colors[i % len(colors)]
        shape = add_shape_rect(slide, left, y, box_width, box_height, color,
                               "", 16, WHITE, True)

        # Label (left side)
        label_box = add_textbox(slide, left + Inches(0.3), y + Inches(0.1),
                                Inches(3.5), Inches(0.65),
                                elem["label"], FONT_HEADING, 18, WHITE, True)

        # Sublabel (right side)
        if "sublabel" in elem:
            sub_box = add_textbox(slide, left + Inches(3.8), y + Inches(0.1),
                                  Inches(5), Inches(0.65),
                                  elem["sublabel"], FONT_BODY, 15, RGBColor(0xDD, 0xDD, 0xFF), False)

    if description:
        add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5),
                    description, FONT_BODY, 14, SUBTLE_TEXT, False)

    set_notes(slide, notes)
    return slide


def make_two_column_slide(prs, title, left_content, right_content, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_COLOR)

    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_COLOR)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0),
                title, FONT_HEADING, 32, ACCENT_COLOR, True)

    add_shape_rect(slide, Inches(0.8), Inches(1.35), Inches(3), Inches(0.06), ACCENT_COLOR)

    # Left column
    add_bullet_slide(slide, left_content, left=Inches(0.8), top=Inches(1.8),
                     width=Inches(5.5), height=Inches(5.0), font_size=16)

    # Right column
    add_bullet_slide(slide, right_content, left=Inches(6.8), top=Inches(1.8),
                     width=Inches(5.5), height=Inches(5.0), font_size=16)

    set_notes(slide, notes)
    return slide


def make_table_slide(prs, title, headers, rows, notes, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_COLOR)

    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_COLOR)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0),
                title, FONT_HEADING, 32, ACCENT_COLOR, True)

    add_shape_rect(slide, Inches(0.8), Inches(1.35), Inches(3), Inches(0.06), ACCENT_COLOR)

    num_cols = len(headers)
    num_rows = len(rows) + 1
    table_width = Inches(11.7)
    table_height = Inches(0.5) * num_rows

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.8), Inches(1.8),
        table_width, table_height
    )
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = FONT_HEADING
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = FONT_BODY
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = TEXT_COLOR
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_ACCENT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    set_notes(slide, notes)
    return slide


def make_code_slide(prs, title, code_text, notes, description=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_COLOR)

    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_COLOR)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0),
                title, FONT_HEADING, 32, ACCENT_COLOR, True)

    add_shape_rect(slide, Inches(0.8), Inches(1.35), Inches(3), Inches(0.06), ACCENT_COLOR)

    # Code background
    code_bg = add_shape_rect(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.5),
                             RGBColor(0x1E, 0x1E, 0x2E))

    # Code text
    add_textbox(slide, Inches(1.1), Inches(1.9), Inches(11.1), Inches(4.1),
                code_text, FONT_MONO, 15, RGBColor(0xD4, 0xD4, 0xD4))

    if description:
        add_textbox(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.8),
                    description, FONT_BODY, 16, TEXT_COLOR)

    set_notes(slide, notes)
    return slide


def make_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_COLOR)

    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), ACCENT_COLOR)

    # Logo
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(5.167), Inches(0.5), Inches(3), Inches(1.875))

    add_textbox(slide, Inches(1), Inches(2.6), Inches(11.333), Inches(0.8),
                "Get Started with PromptKit",
                FONT_HEADING, 36, ACCENT_COLOR, True, PP_ALIGN.CENTER)

    links = [
        "\U0001F517  aka.ms/PromptKit",
        "\U0001F4E6  npx @alan-jowett/promptkit",
        "\U0001F4C4  MIT License — open source",
        "\U0001F419  github.com/microsoft/promptkit",
    ]
    add_bullet_slide(slide, links, left=Inches(3.5), top=Inches(3.7),
                     width=Inches(6.333), height=Inches(2.0), font_size=20,
                     color=TEXT_COLOR, spacing=Pt(12))

    add_textbox(slide, Inches(1), Inches(5.8), Inches(11.333), Inches(0.5),
                "157 components · 15 personas · 48 protocols · 21 formats · 64 templates",
                FONT_BODY, 16, SUBTLE_TEXT, False, PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(6.5), Inches(11.333), Inches(0.8),
                "Thank you!  Questions?",
                FONT_HEADING, 32, ACCENT_COLOR, True, PP_ALIGN.CENTER)

    add_shape_rect(slide, Inches(0), Inches(7.2), SLIDE_WIDTH, Inches(0.3), ACCENT_COLOR)

    set_notes(slide, "PromptKit is open source under MIT license. You can find it at aka.ms/PromptKit or install it with npx. The full repository has 157 components across 6 engineering domains. I encourage you to try it on your next task — whether that's investigating a bug, reviewing code, or writing a requirements document. Thank you! I'm happy to take questions.")
    return slide


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # --- Slide 1: Title ---
    make_title_slide(prs)

    # --- Slide 2: The Hook ---
    make_content_slide(prs,
        "The Most Important Code You're Not Engineering",
        [
            "Every AI-assisted task — investigating bugs, writing requirements,\nreviewing code — lives or dies by the prompt that drives it",
            "",
            "Yet most teams still write these prompts ad hoc:",
            "    • Copy-pasted from chat histories",
            "    • Untested — no way to detect regressions",
            "    • Inconsistent — every engineer writes their own",
            "    • Impossible to improve systematically",
            "",
            "Your code is reviewed, tested, and version-controlled.\nYour prompts? Probably pasted from a Slack thread.",
        ],
        "Think about the last time you used an AI tool — Copilot, ChatGPT, Claude. How much of the result quality came from the prompt you wrote? Probably most of it. Now think about how much engineering rigor you put into that prompt versus the code you wrote afterward. There's a huge gap. We version-control our code. We review it. We test it. But the prompts that drive our AI-assisted engineering? They're throwaway artifacts — copy-pasted, tweaked, and forgotten. Today I want to show you how to close that gap."
    )

    # --- Slide 3: The Problem ---
    make_content_slide(prs,
        "Why Ad-Hoc Prompting Fails at Scale",
        [
            "Non-deterministic — same prompt, different results every time;\nLLMs are inherently stochastic",
            "",
            "Not version-controlled — prompts live in chat histories,\nSlack threads, sticky notes",
            "",
            "Not testable — no way to know if a prompt \"regressed\"\nafter you changed it",
            "",
            "Not reusable — every engineer writes their own from scratch",
            "",
            "Not composable — can't mix and match proven reasoning patterns",
        ],
        "Here's the fundamental problem. LLMs are non-deterministic — you give the same prompt twice, you might get different results. But we make it worse by not applying any engineering discipline. Prompts aren't in source control. We don't test them. We don't share them. When someone on your team discovers a great prompt for code review, how does the rest of the team benefit? Usually, they don't. They write their own version from scratch. And when an LLM hallucinates or misses something critical, we blame the model. But often, the prompt was the problem."
    )

    # --- Slide 4: What If ---
    make_content_slide(prs,
        "What If Prompts Were Engineered Artifacts?",
        [
            "Version-controlled in Git — diffs, blame, history, branching",
            "",
            "Modular — compose proven components instead of writing from scratch",
            "",
            "Testable — compare outputs against known-good reference prompts",
            "",
            "Shareable — team-wide library of best practices, not tribal knowledge",
            "",
            "Improvable — systematic iteration, not random tweaking",
            "",
            "This is what PromptKit does.",
        ],
        "What if we treated prompts the way we treat code? What if they lived in Git, were composed from tested modules, and could be improved systematically? That's the thesis behind PromptKit. It's an open-source library — MIT licensed — that gives you composable, version-controlled prompt components. You don't write prompts from scratch. You compose them from proven building blocks. Let me show you how."
    )

    # --- Slide 5: Section Divider ---
    make_section_divider(prs,
        "Part 1\nHow PromptKit Works",
        "Let's start with the architecture — how PromptKit is structured and how the pieces fit together."
    )

    # --- Slide 6: 5-Layer Architecture ---
    make_diagram_slide(prs,
        "Five Composable Layers",
        [
            {"label": "TEMPLATE", "sublabel": "The task itself — \"Investigate this bug...\""},
            {"label": "FORMAT", "sublabel": "Output structure — investigation report, design doc..."},
            {"label": "TAXONOMY", "sublabel": "Classification — D1–D16 spec drift, K1–K14 kernel defects"},
            {"label": "PROTOCOL", "sublabel": "How to reason — anti-hallucination, root-cause analysis"},
            {"label": "PERSONA", "sublabel": "Who the LLM is — systems engineer, security auditor..."},
        ],
        "PromptKit has five composable layers. At the bottom, a Persona defines who the LLM is — its expertise, its tone, its behavioral constraints. Think of it as the LLM's job title and background. Protocols define how the LLM reasons — systematic methodologies like root-cause analysis, anti-hallucination guardrails, memory safety checks. Formats define what the output looks like — the structure and sections the LLM must produce. Taxonomies provide classification schemes for specific domains. And Templates are the task itself — they compose all the other layers and include task-specific instructions. Each layer is a separate Markdown file that can be mixed and matched.",
        "Each layer is a separate Markdown file — compose by reference, not by copy-paste"
    )

    # --- Slide 7: Personas ---
    make_content_slide(prs,
        "Layer 1 — Personas: Who the LLM Is",
        [
            "15 built-in personas across 6 engineering domains",
            "",
            "Software:  systems-engineer · security-auditor · software-architect",
            "Hardware:  electrical-engineer · embedded-firmware-engineer · rf-engineer · mechanical-engineer",
            "Specialty:  protocol-architect · specification-analyst · reverse-engineer · test-engineer",
            "",
            "Each defines: domain expertise, tone, behavioral constraints, reasoning style",
            "",
            "Or define your own custom persona inline",
        ],
        "The first layer is personas. A persona tells the LLM who it is — not just a name, but deep domain expertise, how it should approach problems, and what tone to use. We have 15 built-in personas covering software engineering, hardware design, embedded firmware, protocol engineering, and more. When you tell an LLM it's a senior systems engineer with deep expertise in memory management and concurrency, it reasons differently than when it's a generic assistant. And these aren't just one-liner descriptions — they're detailed behavioral specifications."
    )

    # --- Slide 8: Protocols ---
    make_content_slide(prs,
        "Layer 2 — Protocols: How the LLM Reasons",
        [
            "48 protocols in three categories:",
            "",
            "Guardrails (5) — Cross-cutting safety",
            "    anti-hallucination · self-verification · adversarial-falsification · minimal-edit-discipline · operational-constraints",
            "",
            "Analysis (17) — Domain-specific checks",
            "    memory-safety-c · thread-safety · security-vulnerability · schematic-compliance · kernel-correctness · and more...",
            "",
            "Reasoning (26) — Systematic methodologies",
            "    root-cause-analysis · requirements-elicitation · traceability-audit · protocol-evolution · change-propagation · and more...",
            "",
            "Each is a multi-phase, step-by-step methodology — protocols compose additively",
        ],
        "Protocols are where the real power is. They define how the LLM reasons — not just 'think step by step,' but detailed, multi-phase methodologies. There are three categories. Guardrails are cross-cutting safety protocols — anti-hallucination prevents fabrication, self-verification makes the LLM check its own work. Analysis protocols are domain-specific — memory safety for C, thread safety, security vulnerabilities, even schematic compliance for hardware. Reasoning protocols are systematic approaches — root-cause analysis with hypothesis generation, requirements elicitation, traceability audits. I'll spotlight two of these in detail shortly."
    )

    # --- Slide 9: Formats & Taxonomies ---
    make_content_slide(prs,
        "Layers 3 & 4 — Formats and Taxonomies",
        [
            "21 output formats — structured document templates:",
            "    investigation-report · requirements-doc · design-doc · validation-plan · structured-patch · presentation-kit · agent-instructions · and more...",
            "",
            "5 classification taxonomies:",
            "    specification-drift (D1–D16): Cross-document divergence",
            "    kernel-defect-categories (K1–K14): OS kernel/driver defects",
            "    stack-lifetime-hazards (H1–H5): Memory escape hazards",
            "    cpp-review-patterns (19 labels): C++ code review findings",
            "    protocol-change-categories (PC1–PC8): Protocol changes",
            "",
            "Fun fact: this presentation was built with PromptKit's author-presentation template + presentation-kit format!",
        ],
        "Formats define what the output looks like. When you ask an LLM to investigate a bug, PromptKit ensures the output has specific sections — findings, root cause analysis, evidence, remediation plan — not just a wall of text. We have 21 formats for everything from investigation reports to design documents to presentation kits. In fact, this very presentation was designed using PromptKit's author-presentation template. Taxonomies provide classification labels. When a protocol finds an issue, it classifies it using a defined scheme — D1 through D16 for specification drift, for example. This makes findings actionable and comparable across reviews."
    )

    # --- Slide 10: Templates (code view) ---
    make_code_slide(prs,
        "Layer 5 — Templates: The Task Itself",
        "---\n"
        "name: investigate-bug\n"
        "persona: systems-engineer\n"
        "protocols:\n"
        "  - guardrails/anti-hallucination\n"
        "  - guardrails/self-verification\n"
        "  - guardrails/operational-constraints\n"
        "  - reasoning/root-cause-analysis\n"
        "taxonomies:\n"
        "  - stack-lifetime-hazards\n"
        "format: investigation-report\n"
        "params:\n"
        "  problem_description: \"What is the bug?\"\n"
        "  code_context: \"Relevant source code\"\n"
        "  environment: \"OS, compiler, runtime\"\n"
        "---",
        "Templates are the orchestration layer. Each template composes the other four layers using YAML frontmatter. This example — investigate-bug — says: use the systems-engineer persona, apply anti-hallucination guardrails, self-verification, and root-cause analysis protocols, and output an investigation report. The parameters are what you provide — the problem description and code context. Everything else is standardized. And because templates declare what they produce and consume, they can be chained into pipelines — the output of one becomes the input of the next.",
        "64 templates — each composes persona + protocols + format via YAML frontmatter.\nTemplates declare input/output contracts → enables chaining into pipelines."
    )

    # --- Slide 11: Composition in Action ---
    make_code_slide(prs,
        "How Components Snap Together",
        "# Identity\n"
        "You are a senior systems engineer with deep expertise\n"
        "in memory management, concurrency, and debugging...\n"
        "\n"
        "# Reasoning Protocols\n"
        "## Protocol: Anti-Hallucination Guardrails\n"
        "Rule 1: Every claim must be KNOWN, INFERRED, or ASSUMED...\n"
        "\n"
        "## Protocol: Root Cause Analysis\n"
        "Phase 1: Characterize the symptom precisely...\n"
        "Phase 2: Generate >= 3 hypotheses before investigating...\n"
        "\n"
        "# Output Format\n"
        "Output MUST contain: Executive Summary, Findings,\n"
        "Root Cause, Remediation, Prevention...\n"
        "\n"
        "# Task\n"
        "Investigate: segfault on startup under high connection load...",
        "Here's what the assembled prompt actually looks like. PromptKit reads each component file and includes its full body text verbatim — no summarization, no abbreviation. The identity section comes from the persona file. The reasoning protocols section includes every phase and sub-step from each protocol. The output format section defines the required structure. And the task section is the template with your parameters filled in. The result is a single, coherent Markdown document that you can paste into any LLM session. The key point: because every component is included verbatim, the LLM gets the full reasoning methodology — not a watered-down summary.",
        "Components are included verbatim — no summarization. One coherent Markdown document."
    )

    # --- Slide 12: Section Divider ---
    make_section_divider(prs,
        "Part 2\nProtocol Spotlights",
        "Now let me zoom in on two protocols to show you the kind of reasoning discipline PromptKit enforces."
    )

    # --- Slide 13: Anti-Hallucination ---
    make_content_slide(prs,
        "Spotlight: Anti-Hallucination Guardrails",
        [
            "Epistemic Labeling — every claim must be categorized:",
            "    KNOWN: Directly stated in the provided context",
            "    INFERRED: Reasonable conclusion with explicit reasoning chain",
            "    ASSUMED: Not established — flagged with [ASSUMPTION]",
            "",
            "Refusal to Fabricate — no invented function names, API signatures, version numbers. Use [UNKNOWN] placeholders instead.",
            "",
            "Uncertainty Disclosure — enumerate multiple interpretations, don't choose one silently",
            "",
            "Source Attribution — cite where information came from (\"per requirements doc, section 3.2\")",
            "",
            "Hard trigger: if ASSUMED content exceeds ~30%, stop and request more context",
        ],
        "The anti-hallucination protocol is probably the most impactful single component in PromptKit. It enforces epistemic honesty. Every claim the LLM makes must be labeled — is this directly known from the context I provided? Is it inferred, and if so, show me the reasoning chain? Or is it an assumption — and if so, flag it explicitly. The protocol also forbids fabrication. If the LLM doesn't know a function name, it writes UNKNOWN instead of making something up. If there are multiple interpretations, it must list them all instead of picking one silently. And there's a hard trigger: if more than 30% of the output is assumptions, the LLM must stop and ask for more context. This single protocol eliminates the majority of hallucination issues."
    )

    # --- Slide 14: Root-Cause Analysis ---
    make_content_slide(prs,
        "Spotlight: Root-Cause Analysis Protocol",
        [
            "Phase 1 — Symptom Characterization",
            "    Describe precisely, establish timeline, determine blast radius",
            "",
            "Phase 2 — Hypothesis Generation",
            "    Generate >= 3 hypotheses before investigating ANY of them",
            "    Include one \"non-obvious\" hypothesis. No anchoring on first guess.",
            "",
            "Phase 3 — Evidence Gathering & Elimination",
            "    For each hypothesis: minimal investigation to confirm/eliminate. Classify: CONFIRMED / ELIMINATED / INCONCLUSIVE",
            "",
            "Phase 4 — Root Cause Identification",
            "    Distinguish root cause from proximate cause. Trace full causal chain. \"If we fix only the proximate cause, will the root cause produce other failures?\"",
            "",
            "Phase 5 — Remediation: Fix the root cause, not the symptom. Propose secondary fixes.",
        ],
        "The root-cause analysis protocol shows what structured reasoning looks like in practice. It's a 5-phase methodology. Phase 1 forces the LLM to precisely characterize the symptom before jumping to conclusions — what's the expected behavior, under what conditions does it fail, what's affected and what isn't. Phase 2 is critical — generate at least three hypotheses before investigating any of them, including at least one non-obvious hypothesis. This prevents anchoring bias. Phase 3 gathers evidence to confirm or eliminate each hypothesis. Phase 4 makes a crucial distinction — the proximate cause versus the root cause. A null pointer dereference is the proximate cause. The missing error handling that left the pointer uninitialized is the root cause. Phase 5 ensures the fix addresses the root cause and proposes preventive measures."
    )

    # --- Slide 15: Section Divider ---
    make_section_divider(prs,
        "Part 3\nWhat Can You Do With PromptKit?",
        "Now that you've seen the building blocks and how the reasoning protocols work, let's talk about the breadth of what you can do with these components."
    )

    # --- Slide 16: Domains Table ---
    make_table_slide(prs,
        "157 Components Across 6 Engineering Domains",
        ["Domain", "What You Can Do"],
        [
            ["Software", "Code review, bug investigation, requirements, design docs, testing, refactoring"],
            ["Hardware / EE", "Schematic review, BOM audit, PCB layout, simulation validation, power budgets"],
            ["Embedded Firmware", "Boot sequences, OTA updates, flash management, power-fail-safe, watchdogs"],
            ["Protocol Engineering", "RFC authoring, protocol evolution, conflict analysis, validation specs"],
            ["Specification Analysis", "Invariant extraction, traceability audits, interface contracts, spec diffing"],
            ["DevOps & CI/CD", "Pipelines, issue triage, PR triage, releases, infrastructure review"],
        ],
        "PromptKit isn't just for software. It covers six engineering domains — software, hardware and electrical engineering, embedded firmware, protocol engineering, specification analysis, and DevOps. 157 components total. You can review C++ code, audit a schematic against requirements, investigate a CI failure, author an RFC, or triage pull requests — all with the same structured, composable approach. And because every template declares what it produces and consumes, they chain together into pipelines.",
        col_widths=[3, 8.7]
    )

    # --- Slide 17: Engineering Lifecycle ---
    make_diagram_slide(prs,
        "The Engineering Lifecycle: Bootstrap → Evolve → Maintain",
        [
            {"label": "BOOTSTRAP", "sublabel": "Scan any repo, extract requirements + design + validation specs"},
            {"label": "EVOLVE", "sublabel": "Propagate changes through specs & implementation with adversarial audits"},
            {"label": "MAINTAIN", "sublabel": "Periodic drift detection — find where code and specs have diverged"},
        ],
        "The three lifecycle workflows are where everything comes together. Bootstrap scans any existing repository and extracts structured specifications — requirements, design, validation — from the code that's already there. You don't need to start from scratch. Evolve takes a requirements change and propagates it through every layer — specs, implementation, tests — with adversarial audits at each transition to catch inconsistencies before they ship. Maintain runs periodically to detect drift — where have the code and specs diverged? — and generates corrective patches. It's a continuous cycle. When maintain detects drift, it loops back to bootstrap.",
        "Each is an interactive workflow — runs in your LLM session with built-in challenge and verification"
    )

    # --- Slide 18: Section Divider ---
    make_section_divider(prs,
        "Part 4\nCase Study — Finding a Crypto\nVulnerability via Spec Review",
        "Let me tell you a real story about what PromptKit found in production."
    )

    # --- Slide 19: Case Study Task ---
    make_content_slide(prs,
        "Case Study: IoT Protocol Crypto Migration",
        [
            "Project: Sonde — a production IoT runtime",
            "    260+ requirements across 15 specification documents",
            "",
            "Task: Replace gateway's asymmetric key-pair pairing encryption with AES-256-GCM using a pre-shared key (PSK)",
            "",
            "Impact: Every radio frame in the protocol is affected",
            "",
            "Approach: PromptKit's engineering-workflow — 8-phase incremental development with adversarial audits at each transition",
            "",
            "Phase 1 discovery: The initial understanding was fundamentally wrong — wrong crypto model, scope significantly underestimated",
        ],
        "Here's a real case study. Sonde is a production IoT runtime with over 260 requirements across 15 specification documents. The task was to replace the gateway's encryption scheme — moving from asymmetric key-pair pairing to AES-256-GCM with a pre-shared key. This affected every radio frame in the protocol. We used PromptKit's engineering workflow — 8 phases, from requirements discovery through implementation with adversarial audits at every transition. Right in Phase 1, the automated analysis discovered that the initial understanding of the crypto model was fundamentally wrong. The scope was far larger than originally estimated."
    )

    # --- Slide 20: The Finding ---
    make_content_slide(prs,
        "Nonce Reuse Vulnerability Found by Automated Spec Review",
        [
            "Round 5 of spec review: automated reviewer identified that GCM nonce construction could allow nonce reuse between request/response pairs using the same PSK",
            "",
            "Original nonce construction:",
            "    gcm_nonce = SHA-256(psk)[0..3] || frame_nonce[8]",
            "",
            "Problem: A request and response could use the same nonce value — catastrophic for GCM security",
            "",
            "Fix: Include msg_type in nonce construction:",
            "    gcm_nonce = SHA-256(psk)[0..3] || msg_type || frame_nonce[8]",
            "",
            "This vulnerability survived 6 rounds of human review. The structured, adversarial prompt caught it.",
        ],
        "In round 5 of the automated specification review, PromptKit's adversarial analysis caught a critical vulnerability. The GCM nonce construction didn't include the message type. This meant a request and a response using the same pre-shared key could end up with the same nonce value. For AES-GCM, nonce reuse is catastrophic — it breaks the authentication guarantees entirely. The fix was straightforward: include the message type in the nonce construction. But here's the key insight — this finding survived six rounds of human review. Domain experts looked at this specification multiple times and didn't catch it. The structured, adversarial prompt did."
    )

    # --- Slide 21: The Lesson ---
    make_content_slide(prs,
        "What the Case Study Tells Us",
        [
            "Full migration: 28 modified, 16 retired, 8 new requirements across 9 specification documents",
            "",
            "PromptKit's contribution:",
            "    • Caught a genuine cryptographic vulnerability that human review missed",
            "    • Forced systematic exploration — no anchoring on first impressions",
            "    • Adversarial auditing at each phase prevented compounding errors",
            "",
            "Human review still mattered:",
            "    • Phase 7 caught a frame ownership error — only domain knowledge could catch it",
            "    • Survived all automated review",
            "",
            "Lesson: Structured prompts + human judgment > either alone",
        ],
        "The full migration resulted in 28 modified, 16 retired, and 8 new requirements across 9 specification documents. PromptKit caught a genuine cryptographic vulnerability that human experts missed repeatedly. But it's not a replacement for human judgment. In Phase 7, user review caught a frame ownership design error that all the automated review missed — the phone generates the PEER_REQUEST ESP-NOW frame, and the node is a pure relay. Only domain knowledge could catch that. The lesson: structured prompts amplify human expertise. They catch what humans miss, and humans catch what automation misses. Together, they're more reliable than either alone."
    )

    # --- Slide 22: Section Divider ---
    make_section_divider(prs,
        "Part 5\nLive Demo",
        "Let's see PromptKit in action."
    )

    # --- Slide 23: Demo ---
    make_content_slide(prs,
        "Demo: Bootstrapping a Prompt with PromptKit",
        [
            "We'll use npx @alan-jowett/promptkit to launch an interactive session",
            "",
            "The bootstrap engine reads manifest.yaml, discovers all 157 components",
            "",
            "We'll select a task, provide parameters, and watch it assemble a structured prompt",
            "",
            "Watch for:",
            "    • Persona selection    • Protocol composition",
            "    • Format enforcement   • Parameter substitution",
        ],
        "[Switch to terminal] Let me show you how this works in practice. I'm going to launch PromptKit and walk through an interactive session. Watch how it selects components, composes them, and produces a structured prompt."
    )

    # --- Slide 23b: Demo Fallback ---
    make_code_slide(prs,
        "Demo Fallback: What the Assembled Prompt Looks Like",
        "# Identity\n"
        "You are a senior systems engineer with deep expertise in\n"
        "memory management, concurrency, performance, and debugging.\n"
        "You reason from first principles...\n"
        "\n"
        "# Reasoning Protocols\n"
        "## Anti-Hallucination Guardrails\n"
        "Every claim MUST be categorized as KNOWN, INFERRED, or ASSUMED.\n"
        "Do NOT invent function names, API signatures...\n"
        "\n"
        "## Root Cause Analysis\n"
        "Phase 1: Describe the symptom precisely...\n"
        "Phase 2: Generate >= 3 hypotheses before investigating...\n"
        "\n"
        "# Output Format\n"
        "Findings, Root Cause, Remediation, Prevention...\n"
        "\n"
        "# Task\n"
        "Investigate: {{problem_description}}",
        "[Use this slide if the demo fails] Here's what an assembled prompt looks like. Notice how every component is included in full — the persona text, every phase of every protocol, the format rules, and the task with parameters filled in. This is a single document you can paste into any LLM session.",
        "Assembled prompt — all components included verbatim, parameters filled in"
    )

    # --- Slide 24: Section Divider ---
    make_section_divider(prs,
        "Part 6\nGetting Started",
        "So how do you actually start using this?"
    )

    # --- Slide 25: Quick Start ---
    make_content_slide(prs,
        "Three Ways to Use PromptKit",
        [
            "1. npx (recommended — no clone needed):",
            "       npx @alan-jowett/promptkit",
            "",
            "2. Clone the repo — full access to all components:",
            "       git clone https://github.com/microsoft/promptkit.git",
            "       copilot -i \"Read and execute bootstrap.md\"",
            "",
            "3. Manual — paste bootstrap.md + manifest.yaml into any LLM session",
            "",
            "Works with: GitHub Copilot CLI · Claude Code · Any LLM",
        ],
        "There are three ways to get started. The easiest is npx — you don't even need to clone the repo. Just run npx @alan-jowett/promptkit and it launches an interactive session. It auto-detects whether you're using GitHub Copilot or Claude Code. If you want full access to all the component files, clone the repo and point your LLM CLI at bootstrap.md. And if you're not using a CLI tool, you can paste the bootstrap prompt and manifest into any LLM session — ChatGPT, Claude, whatever — and follow the interactive flow. It works everywhere."
    )

    # --- Slide 26: Agent Instructions ---
    make_content_slide(prs,
        "Persistent Agent Instructions —\nVersion-Control Your Team's AI Practices",
        [
            "PromptKit produces persistent instruction files for:",
            "",
            "GitHub Copilot: .github/instructions/*.instructions.md — composable skill files with applyTo glob targeting",
            "Claude Code: CLAUDE.md",
            "Cursor: .cursorrules",
            "",
            "Example: \"Always review C code for memory safety\" → persistent skill file applied to every *.c and *.h file",
            "",
            "Commit to your repo → every engineer gets the same AI behavior automatically",
        ],
        "Beyond one-off prompts, PromptKit can produce persistent instruction files. For GitHub Copilot, these are composable skill files under .github/instructions — each with a description and file-targeting via applyTo globs. For Claude Code, it produces a CLAUDE.md file. For Cursor, a .cursorrules file. These go into your repository. Every engineer on your team gets the same AI behavior. If your team decides that C code should always be reviewed for memory safety, that becomes a version-controlled skill file — not tribal knowledge."
    )

    # --- Slide 27: Takeaways ---
    make_content_slide(prs,
        "Key Takeaways",
        [
            "1.  Prompts are code — treat them with the same engineering rigor as your software",
            "",
            "2.  Compose, don't copy-paste — modular components beat ad-hoc prompts every time",
            "",
            "3.  Structured reasoning bounds non-determinism — protocols like anti-hallucination and root-cause analysis make LLM outputs more reliable",
            "",
            "4.  Test against references — compare outputs against known-good prompts to catch regressions",
            "",
            "5.  Humans + structured prompts > either alone — the crypto case study proves it",
        ],
        "Let me leave you with five takeaways. First, prompts are code. Treat them that way. Second, compose from proven components instead of writing from scratch. Third, structured reasoning protocols bound the non-determinism inherent in LLMs — they won't eliminate it, but they make outputs dramatically more reliable. Fourth, test your prompts — PromptKit supports reference comparison testing. And fifth, structured prompts and human judgment are complementary. The nonce-reuse case study proves it — automation caught what humans missed, and humans caught what automation missed."
    )

    # --- Slide 28: Closing ---
    make_closing_slide(prs)

    # Save
    output_path = os.path.join(OUTPUT_DIR, "promptkit-office-hour.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
